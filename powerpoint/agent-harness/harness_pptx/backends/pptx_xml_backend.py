"""PPTX XML backend — cross-platform OOXML manipulation via python-pptx.

Provides a complete cross-platform PowerPoint backend that works on Windows,
Linux, and macOS without requiring Microsoft PowerPoint to be installed.
Uses python-pptx for direct OOXML manipulation.
"""

from __future__ import annotations

from typing import Any

from harness_pptx.backends.interface import (
    ChartStyle,
    ElementHandle,
    PresentationHandle,
    RendererInterface,
    ShapeStyle,
    SlideHandle,
    TableStyle,
    TextMetrics,
    TextStyle,
)


# ---- python-pptx imports (lazy) ------------------------------------------

def _ensure_pptx():
    try:
        import pptx
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.enum.chart import XL_CHART_TYPE
        return pptx
    except ImportError:
        raise RuntimeError(
            "python-pptx is required for PPTXXmlBackend. "
            "Install it with: pip install python-pptx"
        )


# ---- Helpers -------------------------------------------------------------

def _emu(points: float) -> int:
    """Convert points to EMU (1 pt = 12700 EMU)."""
    return int(round(points * 12700))


def _parse_hex(hex_str: str):
    """Parse '#RRGGBB' or '#RGB' into an RGBColor object."""
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        h = "000000"
    from pptx.dml.color import RGBColor
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _hex_to_tuple(hex_str: str) -> tuple[int, int, int]:
    """Parse '#RRGGBB' into (R, G, B) tuple."""
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    if len(h) != 6:
        return (0, 0, 0)
    return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


_SHAPE_TYPE_MAP = {
    "rectangle": "RECTANGLE",
    "rect": "RECTANGLE",
    "rounded_rectangle": "ROUNDED_RECTANGLE",
    "rounded_rect": "ROUNDED_RECTANGLE",
    "oval": "OVAL",
    "circle": "OVAL",
    "ellipse": "OVAL",
    "triangle": "ISOSCELES_TRIANGLE",
    "diamond": "DIAMOND",
    "pentagon": "PENTAGON",
    "hexagon": "HEXAGON",
    "chevron": "CHEVRON",
    "arrow_right": "RIGHT_ARROW",
    "arrow_left": "LEFT_ARROW",
    "star": "STAR_5_POINT",
}

_ALIGN_MAP = {
    "left": "LEFT",
    "center": "CENTER",
    "right": "RIGHT",
    "justify": "JUSTIFY",
}

_CHART_TYPE_MAP = {
    "bar": "BAR_CLUSTERED",
    "bar_stacked": "BAR_STACKED",
    "column": "COLUMN_CLUSTERED",
    "column_stacked": "COLUMN_STACKED",
    "line": "LINE",
    "pie": "PIE",
    "area": "AREA",
    "scatter": "SCATTER",
}


# ---- Backend -------------------------------------------------------------

class PPTXXmlBackend(RendererInterface):
    """Cross-platform backend using python-pptx for OOXML manipulation.

    Works on Windows, Linux, and macOS without Microsoft PowerPoint.
    Native PPTX tables, charts, and shapes via python-pptx.
    """

    backend_name = "pptx-xml"
    platform = "cross"

    def __init__(self):
        self._prs = None
        self._element_counter = 0

    def _next_eid(self, element_type: str = "element") -> str:
        self._element_counter += 1
        return f"elem-{self._element_counter}"

    def _get_slide(self, slide: SlideHandle):
        """Get the python-pptx slide object from a handle."""
        if self._prs is None:
            raise RuntimeError("No presentation open")
        return self._prs.slides[slide.index]

    # ---- Presentation lifecycle ------------------------------------------

    def create_presentation(self) -> PresentationHandle:
        _ensure_pptx()
        from pptx import Presentation
        self._prs = Presentation()
        self._prs.slide_width = _emu(960)
        self._prs.slide_height = _emu(540)
        self._element_counter = 0
        return PresentationHandle(id="pptx-1")

    def open_presentation(self, path: str) -> PresentationHandle:
        _ensure_pptx()
        from pptx import Presentation
        self._prs = Presentation(path)
        self._element_counter = 0
        return PresentationHandle(id="pptx-1", path=path)

    # ---- Slide operations ------------------------------------------------

    def add_slide(self, pres: PresentationHandle, layout: str = "blank") -> SlideHandle:
        if self._prs is None:
            raise RuntimeError("No presentation open")
        slide_layout = self._prs.slide_layouts[6]  # blank layout
        self._prs.slides.add_slide(slide_layout)
        idx = len(self._prs.slides) - 1
        return SlideHandle(id=f"slide-{idx}", index=idx)

    def delete_slide(self, pres: PresentationHandle, slide: SlideHandle) -> None:
        if self._prs is None:
            raise RuntimeError("No presentation open")
        idx = slide.index
        if idx >= len(self._prs.slides):
            return
        try:
            sld_id_lst = self._prs.slides._sldIdLst
            sld_id = sld_id_lst[idx]
            rId = sld_id.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
            )
            if rId is None:
                rId = sld_id.rId
            self._prs.part.drop_rel(rId)
            del sld_id_lst[idx]
        except Exception:
            raise RuntimeError(
                "delete_slide failed — python-pptx does not provide a stable "
                "public API for slide deletion. This may fail across versions."
            )

    def slide_count(self, pres: PresentationHandle) -> int:
        if self._prs is None:
            return 0
        return len(self._prs.slides)

    # ---- Slide background ------------------------------------------------

    def set_slide_background(self, slide: SlideHandle, color: str) -> None:
        from pptx.oxml.ns import qn
        s = self._get_slide(slide)
        bg = s.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = _parse_hex(color)

    # ---- Text ------------------------------------------------------------

    def add_text_box(
        self, slide: SlideHandle,
        x: float, y: float, w: float, h: float,
        text: str, style: TextStyle,
    ) -> ElementHandle:
        from pptx.util import Pt
        s = self._get_slide(slide)
        from pptx.enum.text import PP_ALIGN
        txBox = s.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if text:
            p.text = text
        p.font.size = Pt(style.font_size)
        p.font.color.rgb = _parse_hex(style.font_color)
        p.font.bold = style.bold
        p.font.italic = style.italic
        p.font.name = style.font_name
        p.line_spacing = Pt(style.font_size * style.line_spacing)
        align_name = _ALIGN_MAP.get(style.alignment, "LEFT")
        p.alignment = getattr(PP_ALIGN, align_name)
        eid = self._next_eid("text")
        return ElementHandle(id=eid, element_type="text")

    def set_text(self, slide: SlideHandle, element: ElementHandle, text: str) -> None:
        s = self._get_slide(slide)
        for shape in s.shapes:
            if shape.shape_id == element.metadata.get("shape_id"):
                if shape.has_text_frame:
                    shape.text_frame.paragraphs[0].text = text
                return

    # ---- Image -----------------------------------------------------------

    def add_image(
        self, slide: SlideHandle,
        x: float, y: float, w: float, h: float,
        path: str, fit: str = "contain",
    ) -> ElementHandle:
        s = self._get_slide(slide)
        try:
            pic = s.shapes.add_picture(
                str(path), _emu(x), _emu(y), _emu(w), _emu(h)
            )
        except FileNotFoundError:
            # Fallback: add a placeholder rectangle with the path as text
            from pptx.util import Pt
            txBox = s.shapes.add_textbox(_emu(x), _emu(y), _emu(w), _emu(h))
            txBox.text_frame.paragraphs[0].text = f"[Image: {path}]"
            txBox.text_frame.paragraphs[0].font.size = Pt(10)
            eid = self._next_eid("image_placeholder")
            return ElementHandle(id=eid, element_type="image")
        eid = self._next_eid("image")
        return ElementHandle(
            id=eid, element_type="image",
            metadata={"shape_id": pic.shape_id},
        )

    # ---- Shapes ----------------------------------------------------------

    def add_shape(
        self, slide: SlideHandle, shape_type: str,
        x: float, y: float, w: float, h: float, style: ShapeStyle,
    ) -> ElementHandle:
        from pptx.enum.shapes import MSO_SHAPE
        s = self._get_slide(slide)

        mso_name = _SHAPE_TYPE_MAP.get(shape_type, "RECTANGLE")
        mso_shape = getattr(MSO_SHAPE, mso_name, MSO_SHAPE.RECTANGLE)

        shape = s.shapes.add_shape(mso_shape, _emu(x), _emu(y), _emu(w), _emu(h))

        if style.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = _parse_hex(style.fill_color)
        else:
            shape.fill.background()

        if style.line_color:
            shape.line.color.rgb = _parse_hex(style.line_color)
            from pptx.util import Pt
            shape.line.width = Pt(style.line_weight)
        else:
            shape.line.fill.background()

        # Note: corner radius for rounded rectangles is controlled by
        # MSO_SHAPE.ROUNDED_RECTANGLE itself; fine-tuning the radius
        # requires XML-level adjustment of the adjustValues.

        eid = self._next_eid("shape")
        return ElementHandle(
            id=eid, element_type="shape",
            metadata={"shape_id": shape.shape_id},
        )

    # ---- Lines -----------------------------------------------------------

    def add_line(
        self, slide: SlideHandle,
        x1: float, y1: float, x2: float, y2: float,
        color: str = "#000000", weight: float = 1.0,
    ) -> ElementHandle:
        from pptx.enum.shapes import MSO_CONNECTOR
        s = self._get_slide(slide)
        connector = s.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, _emu(x1), _emu(y1), _emu(x2), _emu(y2)
        )
        connector.line.color.rgb = _parse_hex(color)
        from pptx.util import Pt
        connector.line.width = Pt(weight)
        eid = self._next_eid("line")
        return ElementHandle(
            id=eid, element_type="line",
            metadata={"shape_id": connector.shape_id},
        )

    # ---- Tables ----------------------------------------------------------

    def add_table(
        self, slide: SlideHandle,
        rows: int, cols: int, data: list[list[str]],
        x: float, y: float, w: float, h: float, style: TableStyle,
    ) -> ElementHandle:
        from pptx.util import Pt
        s = self._get_slide(slide)

        table_shape = s.shapes.add_table(rows, cols, _emu(x), _emu(y), _emu(w), _emu(h))
        table = table_shape.table

        # Set column widths evenly
        col_w = _emu(w / cols) if cols > 0 else _emu(w)
        for ci in range(cols):
            table.columns[ci].width = col_w

        for ri in range(rows):
            for ci in range(cols):
                cell = table.cell(ri, ci)
                val = ""
                if ri < len(data) and ci < len(data[ri]):
                    val = str(data[ri][ci])
                cell.text = val

                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.size = Pt(style.font_size)
                    if ri == 0:
                        # Header row
                        paragraph.font.bold = True
                        paragraph.font.color.rgb = _parse_hex(style.header_text)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _parse_hex(style.header_fill)
                    else:
                        paragraph.font.color.rgb = _parse_hex("#000000")
                        if ri % 2 == 0:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _parse_hex(style.body_fill)
                        else:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = _parse_hex(style.alt_body_fill)

        eid = self._next_eid("table")
        return ElementHandle(
            id=eid, element_type="table",
            metadata={"shape_id": table_shape.shape_id},
        )

    # ---- Charts ----------------------------------------------------------

    def add_chart(
        self, slide: SlideHandle, chart_type: str,
        data: dict[str, Any],
        x: float, y: float, w: float, h: float, style: ChartStyle,
    ) -> ElementHandle:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE

        s = self._get_slide(slide)

        chart_data = CategoryChartData()
        categories = data.get("categories", [])
        chart_data.categories = categories

        series_list = data.get("series", [])
        for series in series_list:
            name = series.get("name", "")
            values = series.get("values", [])
            chart_data.add_series(name, values)

        xl_name = _CHART_TYPE_MAP.get(chart_type, "BAR_CLUSTERED")
        xl_type = getattr(XL_CHART_TYPE, xl_name, XL_CHART_TYPE.BAR_CLUSTERED)

        chart_shape = s.shapes.add_chart(
            xl_type, _emu(x), _emu(y), _emu(w), _emu(h), chart_data
        )
        chart = chart_shape.chart

        if not style.show_legend:
            chart.has_legend = False

        if data.get("title"):
            chart.has_title = True
            chart.chart_title.text_frame.paragraphs[0].text = str(data["title"])

        eid = self._next_eid("chart")
        return ElementHandle(
            id=eid, element_type="chart",
            metadata={"shape_id": chart_shape.shape_id},
        )

    # ---- Group -----------------------------------------------------------

    def add_group(
        self, slide: SlideHandle, elements: list[ElementHandle],
    ) -> ElementHandle:
        eid = self._next_eid("group")
        return ElementHandle(
            id=eid, element_type="group",
            metadata={"children": [e.id for e in elements]},
        )

    # ---- Element management ----------------------------------------------

    def delete_element(self, slide: SlideHandle, element: ElementHandle) -> None:
        s = self._get_slide(slide)
        shape_id = element.metadata.get("shape_id") if element.metadata else None
        if shape_id is not None:
            for shape in s.shapes:
                if shape.shape_id == shape_id:
                    sp = shape._element
                    sp.getparent().remove(sp)
                    return

    def set_z_order(self, slide: SlideHandle, element: ElementHandle, action: str) -> None:
        s = self._get_slide(slide)
        shape_id = element.metadata.get("shape_id") if element.metadata else None
        if shape_id is None:
            return
        for shape in s.shapes:
            if shape.shape_id == shape_id:
                sp = shape._element
                parent = sp.getparent()
                if action == "front":
                    parent.append(sp)
                elif action == "back":
                    parent.insert(0, sp)
                return

    # ---- Shape listing ---------------------------------------------------

    def list_shapes(self, slide: SlideHandle) -> list[dict[str, Any]]:
        s = self._get_slide(slide)
        result = []
        for shape in s.shapes:
            info = {
                "shape_id": shape.shape_id,
                "name": shape.name,
                "left": shape.left / 12700 if shape.left else 0,
                "top": shape.top / 12700 if shape.top else 0,
                "width": shape.width / 12700 if shape.width else 0,
                "height": shape.height / 12700 if shape.height else 0,
                "has_text": shape.has_text_frame,
            }
            if shape.has_text_frame:
                info["text"] = shape.text_frame.text
            result.append(info)
        return result

    # ---- Save / Export ---------------------------------------------------

    def save(self, pres: PresentationHandle, path: str) -> None:
        if self._prs is None:
            raise RuntimeError("No presentation open")
        self._prs.save(str(path))

    def export_pdf(self, pres: PresentationHandle, path: str) -> None:
        """Export to PDF.

        python-pptx cannot export to PDF directly. This method attempts to
        use LibreOffice if available, otherwise raises NotImplementedError
        with installation guidance.
        """
        # Try LibreOffice first
        import shutil
        import subprocess
        import tempfile
        import os
        from pathlib import Path

        soffice = shutil.which("soffice")
        if not soffice:
            soffice = shutil.which("libreoffice")

        if soffice and self._prs is not None:
            # Save to temp file, then convert
            fd, tmp_pptx = tempfile.mkstemp(suffix=".pptx")
            os.close(fd)
            try:
                self._prs.save(tmp_pptx)
                out_dir = str(Path(path).parent) or "."
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "pdf",
                     "--outdir", out_dir, tmp_pptx],
                    timeout=60, check=True, capture_output=True,
                )
                # Rename output to desired path
                tmp_pdf = Path(out_dir) / (Path(tmp_pptx).stem + ".pdf")
                if tmp_pdf.exists() and tmp_pdf != Path(path):
                    tmp_pdf.rename(path)
            finally:
                try:
                    os.unlink(tmp_pptx)
                except OSError:
                    pass
            return

        raise NotImplementedError(
            "PDF export requires LibreOffice (soffice) to be installed "
            "and available on PATH. Install it from https://www.libreoffice.org/"
        )

    def export_png(self, pres: PresentationHandle, slide_index: int, path: str) -> None:
        """Export a single slide as PNG.

        Uses LibreOffice if available, otherwise raises NotImplementedError.
        """
        import shutil
        import subprocess
        import tempfile
        import os
        from pathlib import Path

        soffice = shutil.which("soffice")
        if not soffice:
            soffice = shutil.which("libreoffice")

        if soffice and self._prs is not None:
            # Create a temp presentation with just the desired slide
            fd, tmp_pptx = tempfile.mkstemp(suffix=".pptx")
            os.close(fd)
            try:
                self._prs.save(tmp_pptx)
                out_dir = str(Path(path).parent) or "."
                subprocess.run(
                    [soffice, "--headless", "--convert-to", "png",
                     "--outdir", out_dir, tmp_pptx],
                    timeout=60, check=True, capture_output=True,
                )
                # Rename output
                base = Path(tmp_pptx).stem
                src = Path(out_dir) / f"{base}.png"
                if src.exists() and src != Path(path):
                    src.rename(path)
            finally:
                try:
                    os.unlink(tmp_pptx)
                except OSError:
                    pass
            return

        raise NotImplementedError(
            "PNG export requires LibreOffice (soffice) to be installed "
            "and available on PATH."
        )

    # ---- Close -----------------------------------------------------------

    def close(self, pres: PresentationHandle) -> None:
        self._prs = None
        self._element_counter = 0
